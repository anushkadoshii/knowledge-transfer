import os
from typing import List, Dict, Union
from io import BytesIO

from docx import Document
from PyPDF2 import PdfReader
from openpyxl import load_workbook
from pptx import Presentation

def extract_text_from_docx(file_stream: BytesIO) -> str:
    """Extract text from a DOCX file."""
    document = Document(file_stream)
    full_text = []
    for para in document.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)

def extract_text_from_pdf(file_stream: BytesIO) -> str:
    """Extract text from a PDF file."""
    reader = PdfReader(file_stream)
    full_text = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            full_text.append(text)
    return '\n'.join(full_text)

def extract_text_from_xlsx(file_stream: BytesIO) -> str:
    """Extract text from an XLSX file."""
    wb = load_workbook(filename=file_stream, data_only=True)
    full_text = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = [str(cell) for cell in row if cell is not None]
            if row_text:
                full_text.append('\t'.join(row_text))
    return '\n'.join(full_text)

def extract_text_from_pptx(file_stream: BytesIO) -> str:
    """Extract text from a PPTX file."""
    prs = Presentation(file_stream)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return '\n'.join(full_text)

def extract_text_from_txt(file_stream: BytesIO) -> str:
    """Extract text from a TXT file."""
    file_stream.seek(0)
    return file_stream.read().decode('utf-8')

def extract_text_from_file(file_name: str, file_stream: BytesIO) -> str:
    """Detect file type and extract text accordingly."""
    ext = os.path.splitext(file_name)[1].lower()
    if ext == '.pdf':
        return extract_text_from_pdf(file_stream)
    elif ext == '.docx':
        return extract_text_from_docx(file_stream)
    elif ext == '.xlsx':
        return extract_text_from_xlsx(file_stream)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_stream)
    elif ext == '.txt':
        return extract_text_from_txt(file_stream)
    else:
        return ""  # Unsupported file type

def process_uploaded_files(uploaded_files: List) -> List[Dict[str, Union[str, BytesIO]]]:
    """
    Process uploaded files and extract important text content for knowledge transfer.
    Returns a list of dicts: [{'filename': ..., 'content': ...}, ...]
    """
    results = []
    for uploaded_file in uploaded_files:
        file_name = uploaded_file.name
        file_stream = BytesIO(uploaded_file.getvalue())
        try:
            extracted_text = extract_text_from_file(file_name, file_stream)
        except Exception as e:
            extracted_text = f"Error extracting text: {str(e)}"
        results.append({
            "filename": file_name,
            "content": extracted_text
        })
    return results

